"""
Automated Reporter Agent for Autonomous Multi-Agent Business Intelligence System - Phase 6

This agent generates professional PDF reports and PowerPoint presentations
from SQL query results, analytics, and research data.

Features:
- PDF report generation with custom letterheads and formatting
- Executive PowerPoint deck creation (3-slide format)
- Integration with Scientist Agent (charts/stats) and Researcher Agent (market insights)
- Professional templates with branding elements
"""

import os
import io
import base64
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Optional, Any
import json

# PDF generation (optional at import-time)
try:
    from fpdf import FPDF  # provided by fpdf2
except Exception:  # pragma: no cover
    FPDF = None  # type: ignore

# PowerPoint generation (optional at import-time)
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except Exception:  # pragma: no cover
    Presentation = None  # type: ignore
    Inches = None  # type: ignore
    Pt = None  # type: ignore
    PP_ALIGN = None  # type: ignore
    RGBColor = None  # type: ignore

# For embedding charts in reports
import plotly.graph_objects as go
from plotly.io import to_image


class DataGeniePDF(FPDF):
    """Custom PDF class with Autonomous Multi-Agent Business Intelligence System branding"""
    
    def __init__(self, report_title: str = "Autonomous Multi-Agent Business Intelligence System Report", *args, **kwargs):
        super().__init__(*args, **kwargs)
        self.report_title = report_title
        self.company_name = "Autonomous Multi-Agent Business Intelligence System"
        self.page_number = 0
        
    def header(self):
        """Custom header with Autonomous Multi-Agent Business Intelligence System branding"""
        # Logo placeholder (gradient background)
        self.set_fill_color(102, 126, 234)  # Purple gradient color
        self.rect(10, 10, 190, 15, 'F')
        
        # Company name
        self.set_font('Arial', 'B', 16)
        self.set_text_color(255, 255, 255)
        self.set_xy(15, 13)
        self.cell(0, 10, self.company_name, 0, 1, 'L')
        
        # Report title
        self.set_font('Arial', 'I', 10)
        self.set_xy(15, 22)
        self.cell(0, 5, self.report_title, 0, 1, 'L')
        
        # Line break
        self.ln(15)
        
    def footer(self):
        """Custom footer with page numbers"""
        self.set_y(-15)
        self.set_font('Arial', 'I', 8)
        self.set_text_color(128, 128, 128)
        
        # Page number
        self.cell(0, 10, f'Page {self.page_no()}', 0, 0, 'C')
        
        # Generation timestamp
        self.set_x(150)
        self.cell(0, 10, f'Generated: {datetime.now().strftime("%Y-%m-%d %H:%M")}', 0, 0, 'R')
        
    def chapter_title(self, title: str):
        """Add a chapter title"""
        self.set_font('Arial', 'B', 14)
        self.set_text_color(102, 126, 234)
        self.cell(0, 10, title, 0, 1, 'L')
        self.ln(2)
        
    def chapter_body(self, body: str):
        """Add chapter body text"""
        self.set_font('Arial', '', 11)
        self.set_text_color(0, 0, 0)
        self.multi_cell(0, 6, body)
        self.ln()
        
    def add_key_value(self, key: str, value: str):
        """Add a key-value pair with formatting"""
        self.set_font('Arial', 'B', 11)
        self.set_text_color(102, 126, 234)
        self.cell(60, 8, key + ':', 0, 0, 'L')
        
        self.set_font('Arial', '', 11)
        self.set_text_color(0, 0, 0)
        self.cell(0, 8, value, 0, 1, 'L')
        
    def add_table(self, headers: List[str], data: List[List[str]], col_widths: Optional[List[int]] = None):
        """Add a formatted table"""
        if not col_widths:
            col_widths = [190 // len(headers)] * len(headers)
            
        # Header row
        self.set_fill_color(102, 126, 234)
        self.set_text_color(255, 255, 255)
        self.set_font('Arial', 'B', 10)
        
        for i, header in enumerate(headers):
            self.cell(col_widths[i], 8, header, 1, 0, 'C', True)
        self.ln()
        
        # Data rows
        self.set_fill_color(240, 240, 240)
        self.set_text_color(0, 0, 0)
        self.set_font('Arial', '', 10)
        
        for i, row in enumerate(data):
            fill = i % 2 == 0
            for j, cell in enumerate(row):
                self.cell(col_widths[j], 7, str(cell)[:30], 1, 0, 'C', fill)
            self.ln()
        self.ln()


class ReporterAgent:
    """
    Automated Reporter Agent that generates professional reports
    from Autonomous Multi-Agent Business Intelligence System query results, analytics, and research data.
    """
    
    def __init__(self):
        """Initialize the Reporter Agent"""
        self.output_dir = Path("reports")
        self.output_dir.mkdir(exist_ok=True)
        
    def generate_pdf_report(
        self,
        query: str,
        sql_result: Dict[str, Any],
        analytics_result: Optional[Dict[str, Any]] = None,
        research_result: Optional[Dict[str, Any]] = None,
        filename: Optional[str] = None
    ) -> str:
        """
        Generate a comprehensive PDF report
        
        Args:
            query: Original user query
            sql_result: Results from SQL execution
            analytics_result: Results from Scientist Agent (optional)
            research_result: Results from Researcher Agent (optional)
            filename: Custom filename (optional)
            
        Returns:
            Path to generated PDF file
        """
        if FPDF is None:
            raise RuntimeError(
                "PDF generation requires 'fpdf2'. Install with: pip install fpdf2"
            )
        # Generate filename
        if not filename:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"datagenie_report_{timestamp}.pdf"
            
        filepath = self.output_dir / filename
        
        # Create PDF
        pdf = DataGeniePDF(report_title=query[:80])
        pdf.add_page()
        
        # Executive Summary Section
        pdf.chapter_title("1. Executive Summary")
        summary_text = self._generate_executive_summary(query, sql_result, analytics_result, research_result)
        pdf.chapter_body(summary_text)
        
        # Query Details Section
        pdf.chapter_title("2. Query Details")
        pdf.add_key_value("User Query", query)
        pdf.add_key_value("Execution Date", datetime.now().strftime("%Y-%m-%d %H:%M:%S"))
        pdf.add_key_value("Query Type", sql_result.get("method", "Standard"))
        pdf.ln()
        
        # SQL Query Section
        if sql_result.get("sql"):
            pdf.chapter_title("3. Generated SQL Query")
            pdf.set_font('Courier', '', 9)
            pdf.set_text_color(0, 0, 0)
            pdf.multi_cell(0, 5, sql_result["sql"][:500])
            pdf.ln()
        
        # Data Results Section
        if sql_result.get("data"):
            pdf.chapter_title("4. Query Results")
            data = sql_result["data"]
            
            if isinstance(data, list) and len(data) > 0:
                # Display as table (first 10 rows)
                if isinstance(data[0], dict):
                    headers = list(data[0].keys())[:5]  # Limit to 5 columns
                    rows = [[str(row.get(h, ""))[:30] for h in headers] for row in data[:10]]
                    pdf.add_table(headers, rows)
                    
                    if len(data) > 10:
                        pdf.set_font('Arial', 'I', 10)
                        pdf.cell(0, 6, f"Showing 10 of {len(data)} total rows", 0, 1, 'L')
            else:
                pdf.chapter_body(str(data)[:500])
        
        # Analytics Section
        if analytics_result:
            pdf.add_page()
            pdf.chapter_title("5. Statistical Analysis")
            
            # Analysis summary
            if analytics_result.get("analysis"):
                pdf.chapter_body(analytics_result["analysis"][:800])
            
            # Key metrics
            if analytics_result.get("statistics"):
                pdf.ln()
                pdf.set_font('Arial', 'B', 12)
                pdf.cell(0, 8, "Key Statistics:", 0, 1, 'L')
                
                stats = analytics_result["statistics"]
                for key, value in list(stats.items())[:10]:
                    pdf.add_key_value(key.replace("_", " ").title(), str(value))
        
        # Research Section
        if research_result:
            pdf.add_page()
            pdf.chapter_title("6. Market Research Insights")
            
            # Internal findings
            if research_result.get("internal_findings"):
                pdf.set_font('Arial', 'B', 11)
                pdf.cell(0, 8, "Internal Performance:", 0, 1, 'L')
                pdf.chapter_body(research_result["internal_findings"][:500])
            
            # External research
            if research_result.get("external_research"):
                pdf.set_font('Arial', 'B', 11)
                pdf.cell(0, 8, "Market Context:", 0, 1, 'L')
                pdf.chapter_body(research_result["external_research"][:500])
            
            # Unified insights
            if research_result.get("unified_insights"):
                pdf.ln()
                pdf.set_font('Arial', 'B', 11)
                pdf.cell(0, 8, "Strategic Recommendations:", 0, 1, 'L')
                pdf.chapter_body(research_result["unified_insights"][:500])
        
        # Recommendations Section
        pdf.add_page()
        pdf.chapter_title("7. Recommendations & Next Steps")
        recommendations = self._generate_recommendations(sql_result, analytics_result, research_result)
        pdf.chapter_body(recommendations)
        
        # Save PDF
        pdf.output(str(filepath))
        
        return str(filepath)
    
    def generate_pptx_report(
        self,
        query: str,
        sql_result: Dict[str, Any],
        analytics_result: Optional[Dict[str, Any]] = None,
        research_result: Optional[Dict[str, Any]] = None,
        filename: Optional[str] = None
    ) -> str:
        """
        Generate an Executive PowerPoint deck (3 slides)
        
        Args:
            query: Original user query
            sql_result: Results from SQL execution
            analytics_result: Results from Scientist Agent (optional)
            research_result: Results from Researcher Agent (optional)
            filename: Custom filename (optional)
            
        Returns:
            Path to generated PPTX file
        """
        if Presentation is None:
            raise RuntimeError(
                "PowerPoint generation requires 'python-pptx'. Install with: pip install python-pptx"
            )
        # Generate filename
        if not filename:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"datagenie_executive_deck_{timestamp}.pptx"
            
        filepath = self.output_dir / filename
        
        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Slide 1: Overview
        self._add_overview_slide(prs, query, sql_result)
        
        # Slide 2: Key Findings
        self._add_findings_slide(prs, sql_result, analytics_result)
        
        # Slide 3: Market Context
        self._add_market_context_slide(prs, research_result)
        
        # Save presentation
        prs.save(str(filepath))
        
        return str(filepath)
    
    def _add_overview_slide(self, prs: Presentation, query: str, sql_result: Dict[str, Any]):
        """Add overview slide to PowerPoint"""
        # Title slide layout
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        
        # Title
        title = slide.shapes.title
        title.text = "Autonomous Multi-Agent Business Intelligence System Executive Report"
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 126, 234)
        
        # Subtitle
        subtitle = slide.placeholders[1]
        subtitle.text = f"{query[:100]}\n\n{datetime.now().strftime('%B %d, %Y')}"
        subtitle.text_frame.paragraphs[0].font.size = Pt(24)
        
    def _add_findings_slide(self, prs: Presentation, sql_result: Dict[str, Any], analytics_result: Optional[Dict[str, Any]]):
        """Add key findings slide to PowerPoint"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        # Title
        title = slide.shapes.title
        title.text = "Key Findings"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.bold = True
        
        # Content
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()
        
        # Finding 1: Query results
        p = tf.paragraphs[0]
        p.text = "Data Analysis Results"
        p.font.size = Pt(20)
        p.font.bold = True
        p.space_after = Pt(10)
        
        if sql_result.get("data"):
            data = sql_result["data"]
            p = tf.add_paragraph()
            if isinstance(data, list):
                p.text = f"• Retrieved {len(data)} records from database"
            else:
                p.text = f"• Query executed successfully"
            p.font.size = Pt(18)
            p.space_after = Pt(10)
        
        # Finding 2: Analytics insights
        if analytics_result:
            p = tf.add_paragraph()
            p.text = "\nStatistical Insights"
            p.font.size = Pt(20)
            p.font.bold = True
            p.space_after = Pt(10)
            
            if analytics_result.get("analysis"):
                analysis = analytics_result["analysis"][:150]
                p = tf.add_paragraph()
                p.text = f"• {analysis}"
                p.font.size = Pt(18)
                p.space_after = Pt(10)
            
            # Key statistics
            if analytics_result.get("statistics"):
                stats = analytics_result["statistics"]
                for key, value in list(stats.items())[:3]:
                    p = tf.add_paragraph()
                    p.text = f"• {key.replace('_', ' ').title()}: {value}"
                    p.font.size = Pt(18)
                    p.space_after = Pt(5)
    
    def _add_market_context_slide(self, prs: Presentation, research_result: Optional[Dict[str, Any]]):
        """Add market context slide to PowerPoint"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        # Title
        title = slide.shapes.title
        title.text = "Market Context & Recommendations"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.bold = True
        
        # Content
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()
        
        if research_result:
            # Internal performance
            if research_result.get("internal_findings"):
                p = tf.paragraphs[0]
                p.text = "Internal Performance"
                p.font.size = Pt(20)
                p.font.bold = True
                p.space_after = Pt(10)
                
                p = tf.add_paragraph()
                p.text = f"• {research_result['internal_findings'][:120]}"
                p.font.size = Pt(18)
                p.space_after = Pt(15)
            
            # Market context
            if research_result.get("external_research"):
                p = tf.add_paragraph()
                p.text = "Market Benchmarking"
                p.font.size = Pt(20)
                p.font.bold = True
                p.space_after = Pt(10)
                
                p = tf.add_paragraph()
                p.text = f"• {research_result['external_research'][:120]}"
                p.font.size = Pt(18)
                p.space_after = Pt(15)
            
            # Recommendations
            if research_result.get("unified_insights"):
                p = tf.add_paragraph()
                p.text = "Strategic Recommendations"
                p.font.size = Pt(20)
                p.font.bold = True
                p.space_after = Pt(10)
                
                p = tf.add_paragraph()
                p.text = f"• {research_result['unified_insights'][:120]}"
                p.font.size = Pt(18)
        else:
            p = tf.paragraphs[0]
            p.text = "Recommendations"
            p.font.size = Pt(20)
            p.font.bold = True
            p.space_after = Pt(10)
            
            p = tf.add_paragraph()
            p.text = "• Continue monitoring key metrics"
            p.font.size = Pt(18)
            p.space_after = Pt(10)
            
            p = tf.add_paragraph()
            p.text = "• Consider running analytics for deeper insights"
            p.font.size = Pt(18)
            p.space_after = Pt(10)
            
            p = tf.add_paragraph()
            p.text = "• Benchmark against market data for competitive analysis"
            p.font.size = Pt(18)
    
    def _generate_executive_summary(
        self,
        query: str,
        sql_result: Dict[str, Any],
        analytics_result: Optional[Dict[str, Any]],
        research_result: Optional[Dict[str, Any]]
    ) -> str:
        """Generate executive summary text"""
        summary_parts = []
        
        # Introduction
        summary_parts.append(
            f"This report presents the results of the data analysis query: '{query}'. "
            f"The analysis was conducted on {datetime.now().strftime('%B %d, %Y')}."
        )
        
        # Data summary
        if sql_result.get("data"):
            data = sql_result["data"]
            if isinstance(data, list):
                summary_parts.append(f"\n\nThe query retrieved {len(data)} records from the database.")
        
        # Analytics summary
        if analytics_result and analytics_result.get("analysis"):
            summary_parts.append(
                f"\n\nStatistical analysis reveals: {analytics_result['analysis'][:200]}..."
            )
        
        # Research summary
        if research_result and research_result.get("unified_insights"):
            summary_parts.append(
                f"\n\nMarket research indicates: {research_result['unified_insights'][:200]}..."
            )
        
        return "".join(summary_parts)
    
    def _generate_recommendations(
        self,
        sql_result: Dict[str, Any],
        analytics_result: Optional[Dict[str, Any]],
        research_result: Optional[Dict[str, Any]]
    ) -> str:
        """Generate recommendations based on results"""
        recommendations = []
        
        recommendations.append("Based on the analysis, we recommend the following actions:\n")
        
        # Default recommendations
        recommendations.append("\n1. Monitor Trends: Continue tracking the key metrics identified in this report.")
        recommendations.append("\n2. Deep Dive Analysis: Consider running more detailed analytics on specific data segments.")
        recommendations.append("\n3. Competitive Benchmarking: Compare results against industry standards and competitors.")
        recommendations.append("\n4. Action Planning: Develop specific action plans based on the insights discovered.")
        
        # Analytics-specific recommendations
        if analytics_result:
            recommendations.append("\n5. Statistical Validation: Validate statistical findings with additional data sources.")
        
        # Research-specific recommendations
        if research_result:
            recommendations.append("\n6. Market Positioning: Adjust strategy based on market insights.")
        
        recommendations.append("\n\nFor detailed implementation guidance, please consult with your data team.")
        
        return "".join(recommendations)
    
    def generate_combined_report(
        self,
        query: str,
        sql_result: Dict[str, Any],
        analytics_result: Optional[Dict[str, Any]] = None,
        research_result: Optional[Dict[str, Any]] = None,
        formats: List[str] = ["pdf", "pptx"]
    ) -> Dict[str, str]:
        """
        Generate reports in multiple formats
        
        Args:
            query: Original user query
            sql_result: Results from SQL execution
            analytics_result: Results from Scientist Agent (optional)
            research_result: Results from Researcher Agent (optional)
            formats: List of formats to generate (pdf, pptx)
            
        Returns:
            Dictionary mapping format to filepath
        """
        results = {}
        
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        
        if "pdf" in formats:
            pdf_filename = f"datagenie_report_{timestamp}.pdf"
            pdf_path = self.generate_pdf_report(
                query, sql_result, analytics_result, research_result, pdf_filename
            )
            results["pdf"] = pdf_path
        
        if "pptx" in formats:
            pptx_filename = f"datagenie_deck_{timestamp}.pptx"
            pptx_path = self.generate_pptx_report(
                query, sql_result, analytics_result, research_result, pptx_filename
            )
            results["pptx"] = pptx_path
        
        return results


# Example usage and testing
if __name__ == "__main__":
    # Create reporter agent
    reporter = ReporterAgent()
    
    # Sample data
    sample_query = "Show me revenue trends for Q4 2025"
    sample_sql_result = {
        "method": "analytics",
        "sql": "SELECT date, SUM(revenue) as total_revenue FROM sales WHERE date >= '2025-10-01' GROUP BY date",
        "data": [
            {"date": "2025-10-01", "total_revenue": 150000},
            {"date": "2025-11-01", "total_revenue": 175000},
            {"date": "2025-12-01", "total_revenue": 200000}
        ]
    }
    sample_analytics = {
        "analysis": "Revenue shows strong upward trend with 33% growth over the quarter.",
        "statistics": {
            "mean_revenue": 175000,
            "growth_rate": 0.33,
            "total_revenue": 525000
        }
    }
    sample_research = {
        "internal_findings": "Our Q4 revenue grew 33% to $525K total.",
        "external_research": "Industry average Q4 growth was 25% according to market reports.",
        "unified_insights": "We outperformed the industry average by 8 percentage points."
    }
    
    # Generate reports
    print("Generating reports...")
    reports = reporter.generate_combined_report(
        sample_query,
        sample_sql_result,
        sample_analytics,
        sample_research
    )
    
    print(f"PDF Report: {reports.get('pdf')}")
    print(f"PPTX Deck: {reports.get('pptx')}")
